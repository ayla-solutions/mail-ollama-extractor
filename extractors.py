import base64, io, re, os
from typing import List, Tuple
from bs4 import BeautifulSoup
from email import policy
from email.parser import BytesParser
from PIL import Image
import pytesseract

from pdfminer.high_level import extract_text as pdf_extract_text
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import pandas as pd
import msg_parser

ENABLE_OCR = os.getenv("ENABLE_OCR", "true").lower() == "true"

def _from_html(html: str) -> str:
    soup = BeautifulSoup(html, "lxml")
    return soup.get_text(" ", strip=True)

def _from_pdf_bytes(b: bytes) -> str:
    try:
        with io.BytesIO(b) as f:
            return pdf_extract_text(f) or ""
    except Exception:
        # Fallback to PyPDF2
        with io.BytesIO(b) as f:
            reader = PdfReader(f)
            pages = [p.extract_text() or "" for p in reader.pages]
            return "\n".join(pages)

def _from_image_bytes(b: bytes) -> str:
    if not ENABLE_OCR:
        return ""
    img = Image.open(io.BytesIO(b))
    return pytesseract.image_to_string(img)

def _from_docx_bytes(b: bytes) -> str:
    f = io.BytesIO(b)
    d = docx.Document(f)
    return "\n".join([p.text for p in d.paragraphs])

def _from_pptx_bytes(b: bytes) -> str:
    prs = Presentation(io.BytesIO(b))
    txt = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt.append(shape.text)
    return "\n".join(txt)

def _from_excel_bytes(b: bytes) -> str:
    with io.BytesIO(b) as f:
        try:
            df = pd.read_excel(f, sheet_name=None, dtype=str)
            return "\n".join([d.to_string(index=False) for _, d in df.items()])
        except Exception:
            f.seek(0)
            df = pd.read_csv(f, dtype=str, encoding_errors="ignore")
            return df.to_string(index=False)

def _from_eml_bytes(b: bytes) -> str:
    em = BytesParser(policy=policy.default).parsebytes(b)
    body = em.get_body(preferencelist=('html','plain'))
    body_text = _from_html(body.get_content()) if body and body.get_content_type()=="text/html" else (body.get_content() if body else "")
    return f"Subject: {em['subject']}\nFrom: {em['from']}\nTo: {em['to']}\n\n{body_text}"

def _from_msg_bytes(b: bytes) -> str:
    # uses msg-parser which works without MAPI on Linux
    with io.BytesIO(b) as f:
        m = msg_parser.MsgParser(f).parse()
    body = m.body if m.body else ""
    if m.body_html:
        body = _from_html(m.body_html)
    return f"Subject: {m.subject}\nFrom: {m.sender}\nTo: {','.join(m.to or [])}\n\n{body}"

def detect_type(name: str) -> str:
    n = name.lower()
    for ext, t in [
        (".pdf","pdf"),(".png","image"),(".jpg","image"),(".jpeg","image"),(".webp","image"),
        (".docx","docx"),(".pptx","pptx"),(".xlsx","excel"),(".xls","excel"),(".csv","excel"),
        (".txt","txt"),(".html","html"),(".eml","eml"),(".msg","msg")
    ]:
        if n.endswith(ext): return t
    return "bin"

def extract_attachment_text(attachments: List[Tuple[str,str]]) -> Tuple[str,str]:
    """
    attachments: list of (filename, base64string)
    returns (combined_text, names_csv)
    """
    blobs = []
    names = []
    for name, b64 in attachments:
        names.append(name)
        b = base64.b64decode(b64)
        kind = detect_type(name)
        try:
            if kind=="pdf": blobs.append(_from_pdf_bytes(b))
            elif kind=="image": blobs.append(_from_image_bytes(b))
            elif kind=="docx": blobs.append(_from_docx_bytes(b))
            elif kind=="pptx": blobs.append(_from_pptx_bytes(b))
            elif kind=="excel": blobs.append(_from_excel_bytes(b))
            elif kind=="txt": blobs.append(b.decode("utf-8", "ignore"))
            elif kind=="html": blobs.append(_from_html(b.decode("utf-8","ignore")))
            elif kind=="eml": blobs.append(_from_eml_bytes(b))
            elif kind=="msg": blobs.append(_from_msg_bytes(b))
            else: blobs.append("")  # unknown/binary
        except Exception:
            blobs.append("")
    combined = "\n\n".join([t for t in blobs if t])
    return combined, ", ".join(names)
